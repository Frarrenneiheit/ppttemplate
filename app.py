from flask import Flask, render_template, request, send_file
import pandas as pd
from pptx import Presentation
import io
import logging
import re

app = Flask(__name__)
logging.basicConfig(level=logging.DEBUG)

MAPPINGS = {
    'S7PI': ('UnifiedData', 'AD10'),
    'S7PF': ('UnifiedData', 'AD11'),
    'S7PL': ('UnifiedData', 'AD12'),
    'S7PY': ('UnifiedData', 'AD13'),
    'S10C': ('UnifiedData', 'S49'),
    'S10D': ('UnifiedData', 'T50'),
    'S10ImpPP1': ('UnifiedData', 'T34'),
    'S10ImpPP2': ('UnifiedData', 'T35'),
    'S10ImpPP3': ('UnifiedData', 'T36'),
    'S10ImpPP4': ('UnifiedData', 'T37'),
    'S10ImpSc1': ('UnifiedData', 'Y34'),
    'S10ImpSc2': ('UnifiedData', 'Y35'),
    'S10ImpSc3': ('UnifiedData', 'Y36'),
    'S10ImpSc4': ('UnifiedData', 'Y37'),
    'S10ImpBe1': ('Benchmark', 'C5'),
    'S10ImpBe2': ('Benchmark', 'C6'),
    'S10ImpBe3': ('Benchmark', 'C7'),
    'S10ImpBe4': ('Benchmark', 'C8'),
    'S10EngPP1': ('UnifiedData', 'U34'),
    'S10EngPP2': ('UnifiedData', 'U35'),
    'S10EngPP3': ('UnifiedData', 'U36'),
    'S10EngSc1': ('UnifiedData', 'Z34'),
    'S10EngSc2': ('UnifiedData', 'Z35'),
    'S10EngSc3': ('UnifiedData', 'Z36'),
    'S10EngBe1': ('Benchmark', 'D5'),
    'S10EngBe2': ('Benchmark', 'D6'),
    'S10EngBe3': ('Benchmark', 'D7'),
    'S10ERPP1': ('UnifiedData', 'V34'),
    'S10ERPP2': ('UnifiedData', 'V35'),
    'S10ERPP3': ('UnifiedData', 'V36'),
    'S10ERSc1': ('UnifiedData', 'AA34'),
    'S10ERSc2': ('UnifiedData', 'AA35'),
    'S10ERSc3': ('UnifiedData', 'AA36'),
    'S10ERBe1': ('Benchmark', 'E5'),
    'S10ERBe2': ('Benchmark', 'E6'),
    'S10ERBe3': ('Benchmark', 'E7'),
    'S10VV1': ('UnifiedData', 'W36'),
    'S10VV2': ('UnifiedData', 'AB36'),
    'S10VV3': ('Benchmark', 'F7'),
    'S10VR1': ('UnifiedData', 'X36'),
    'S10VR2': ('UnifiedData', 'AC36'),
    'S10VR3': ('Benchmark', 'G7'),
    'S12PF1': ('12 Months Averages', 'S3'),
    'S12PF3': ('UnifiedData', 'S51'),
    'S12PF4': ('UnifiedData', 'T52'),
    'S12ImpPP1': ('UnifiedData', 'T38'),
    'S12ImpPP2': ('UnifiedData', 'T39'),
    'S12ImpPP3': ('UnifiedData', 'T40'),
    'S12ImpSc1': ('UnifiedData', 'Y38'),
    'S12ImpSc2': ('UnifiedData', 'Y39'),
    'S12ImpSc3': ('UnifiedData', 'Y40'),
    'S12ImpBe1': ('Benchmark', 'C2'),
    'S12ImpBe2': ('Benchmark', 'C3'),
    'S12ImpBe3': ('Benchmark', 'C4'),
    'S12EngPP1': ('UnifiedData', 'U38'),
    'S12EngPP2': ('UnifiedData', 'U39'),
    'S12EngPP3': ('UnifiedData', 'U40'),
    'S12EngSc1': ('UnifiedData', 'Z38'),
    'S12EngSc2': ('UnifiedData', 'Z39'),
    'S12EngSc3': ('UnifiedData', 'Z40'),
    'S12EngBe1': ('Benchmark', 'D2'),
    'S12EngBe2': ('Benchmark', 'D3'),
    'S12EngBe3': ('Benchmark', 'D4'),
    'S12ERPP1': ('UnifiedData', 'V38'),
    'S12ERPP2': ('UnifiedData', 'V39'),
    'S12ERPP3': ('UnifiedData', 'V40'),
    'S12ERSc1': ('UnifiedData', 'AA38'),
    'S12ERSc2': ('UnifiedData', 'AA39'),
    'S12ERSc3': ('UnifiedData', 'AA40'),
    'S12ERBe1': ('Benchmark', 'E2'),
    'S12ERBe2': ('Benchmark', 'E3'),
    'S12ERBe3': ('Benchmark', 'E4'),
    'S12VV1': ('Chart Tables Platform', 'F30'),
    'S12VV2': ('Chart Tables Platform', 'K30'),
    'S12VV3': ('Benchmark', 'F4'),
    'S12VR1': ('Chart Tables Platform', 'G30'),
    'S12VR2': ('Chart Tables Platform', 'L30'),
    'S12VR3': ('Benchmark', 'G4'),
    'S14PF3': ('UnifiedData', 'S53'),
    'S14PF4': ('UnifiedData', 'T54'),
    'S14ImpPP1': ('Chart Tables Platform', 'C35'),
    'S14ImpPP2': ('Chart Tables Platform', 'C36'),
    'S14ImpPP3': ('Chart Tables Platform', 'C37'),
    'S14ImpSc1': ('Chart Tables Platform', 'H35'),
    'S14ImpSc2': ('Chart Tables Platform', 'H36'),
    'S14ImpSc3': ('Chart Tables Platform', 'H37'),
    'S14ImpBe1': ('Benchmark', 'C9'),
    'S14ImpBe2': ('Benchmark', 'C10'),
    'S14ImpBe3': ('Benchmark', 'C11'),
    'S14EngPP1': ('Chart Tables Platform', 'D35'),
    'S14EngPP2': ('Chart Tables Platform', 'D36'),
    'S14EngPP3': ('Chart Tables Platform', 'D37'),
    'S14EngSc1': ('Chart Tables Platform', 'I35'),
    'S14EngSc2': ('Chart Tables Platform', 'I36'),
    'S14EngSc3': ('Chart Tables Platform', 'I37'),
    'S14EngBe1': ('Benchmark', 'D9'),
    'S14EngBe2': ('Benchmark', 'D10'),
    'S14EngBe3': ('Benchmark', 'D11'),
    'S14ERPP1': ('Chart Tables Platform', 'E35'),
    'S14ERPP2': ('Chart Tables Platform', 'E36'),
    'S14ERPP3': ('Chart Tables Platform', 'E37'),
    'S14ERSc1': ('Chart Tables Platform', 'J35'),
    'S14ERSc2': ('Chart Tables Platform', 'J36'),
    'S14ERSc3': ('Chart Tables Platform', 'J37'),
    'S14ERBe1': ('Benchmark', 'E9'),
    'S14ERBe2': ('Benchmark', 'E10'),
    'S14ERBe3': ('Benchmark', 'E11'),
    'S14VV1': ('Chart Tables Platform', 'F37'),
    'S14VV2': ('Chart Tables Platform', 'K37'),
    'S14VV3': ('Benchmark', 'F11'),
    'S14VR1': ('Chart Tables Platform', 'G37'),
    'S14VR2': ('Chart Tables Platform', 'L37'),
    'S14VR3': ('Benchmark', 'G11'),
    'S16PF3': ('UnifiedData', 'S55'),
    'S16PF4': ('UnifiedData', 'T56'),
    'S16ImpPP1': ('Chart Tables Platform', 'C42'),
    'S16ImpSc1': ('Chart Tables Platform', 'H42'),
    'S16ImpBe1': ('Benchmark', 'C12'),
    'S16ImpBe2': ('Benchmark', 'C13'),
    'S16EngPP1': ('Chart Tables Platform', 'D42'),
    'S16EngSc1': ('Chart Tables Platform', 'I42'),
    'S16EngBe1': ('Benchmark', 'D12'),
    'S16ERPP1': ('Chart Tables Platform', 'E42'),
    'S16ERSc1': ('Chart Tables Platform', 'J42'),
    'S16ERBe1': ('Benchmark', 'E12'),
    'S16ERBe2': ('Benchmark', 'E13'),
    'S16VV1': ('Chart Tables Platform', 'F42'),
    'S16VV2': ('Chart Tables Platform', 'K42'),
    'S16VV3': ('Benchmark', 'F12'),
    'S16VV6': ('Benchmark', 'F13'),
    'S16VR1': ('Chart Tables Platform', 'G42'),
    'S16VR2': ('Chart Tables Platform', 'L42'),
    'S16VR3': ('Benchmark', 'G12'),
    'S16VR6': ('Benchmark', 'G13'),
}

def excel_cell_to_indices(cell_ref):
    match = re.match(r'([A-Z]+)(\d+)', cell_ref)
    if not match:
        raise ValueError(f"Ungültige Zellenreferenz: {cell_ref}")
    
    col_str, row_str = match.groups()
    
    col_idx = 0
    for i, char in enumerate(reversed(col_str)):
        col_idx += (ord(char.upper()) - ord('A') + 1) * (26 ** i)
    col_idx -= 1
    
    row_idx = int(row_str) - 2
    return row_idx, col_idx

def format_number(value):
    try:
        num = float(value)
        if isinstance(num, float) and num < 1:  # Prozentformat für Zahlen < 1
            return f"{num * 100:.1f}%"
        if num >= 1000:
            if num >= 1000000:
                return f"{num/1000000:.0f}M"
            return f"{num/1000:.0f}k"
        return str(num)
    except ValueError:
        return value

def search_and_replace_in_shape(shape, mappings, excel_data):
    for placeholder, (sheet_name, cell_ref) in mappings.items():
        run_found = False
        
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        try:
                            df = excel_data[sheet_name]
                            row_idx, col_idx = excel_cell_to_indices(cell_ref)
                            value = format_number(df.iloc[row_idx, col_idx])
                            run.text = run.text.replace(placeholder, value)
                            run_found = True
                        except Exception as e:
                            logging.error(f"Fehler beim Ersetzen von {placeholder}: {str(e)}")
        
        # Fallback für andere Shape-Typen wenn kein Run gefunden wurde
        if not run_found:
            if hasattr(shape, "text") and placeholder in shape.text:
                text = shape.text
            elif hasattr(shape, "cell") and hasattr(shape.cell, "text") and placeholder in shape.cell.text:
                text = shape.cell.text
            elif hasattr(shape, "title") and placeholder in shape.title:
                text = shape.title
            else:
                continue
                
            try:
                df = excel_data[sheet_name]
                row_idx, col_idx = excel_cell_to_indices(cell_ref)
                value = format_number(df.iloc[row_idx, col_idx])
                
                if hasattr(shape, "text"):
                    shape.text = text.replace(placeholder, value)
                elif hasattr(shape, "cell"):
                    shape.cell.text = text.replace(placeholder, value)
                elif hasattr(shape, "title"):
                    shape.title = text.replace(placeholder, value)
                    
            except Exception as e:
                logging.error(f"Fehler beim Ersetzen von {placeholder}: {str(e)}")
                
def process_files(excel_file, ppt_template):
    try:
        logging.info("Starting file processing")
        excel_data = pd.read_excel(excel_file, sheet_name=None)
        
        if not excel_data:
            raise ValueError("Excel-Datei ist leer")
            
        prs = Presentation(ppt_template)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                search_and_replace_in_shape(shape, MAPPINGS, excel_data)
                
                try:
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                search_and_replace_in_shape(cell, MAPPINGS, excel_data)
                except (AttributeError, TypeError):
                    pass
                
                try:
                    if hasattr(shape, "shapes"):
                        for subshape in shape.shapes:
                            search_and_replace_in_shape(subshape, MAPPINGS, excel_data)
                except (AttributeError, TypeError):
                    pass
        
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io
        
    except Exception as e:
        raise Exception(f"Fehler bei der Verarbeitung: {str(e)}")
        
@app.route('/', methods=['GET'])
def index():
    return render_template('upload.html')

@app.route('/upload', methods=['POST'])
def upload():
    try:
        if 'excel_file' not in request.files or 'ppt_template' not in request.files:
            return 'Keine Dateien hochgeladen', 400
            
        excel_file = request.files['excel_file']
        ppt_template = request.files['ppt_template']
        
        if excel_file.filename == '' or ppt_template.filename == '':
            return 'Keine Dateien ausgewählt', 400
            
        result = process_files(excel_file, ppt_template)
        
        return send_file(
            result,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='result.pptx'
        )
    except Exception as e:
        return str(e), 400

if __name__ == '__main__':
    app.run(debug=True)